#!/usr/bin/env python3
"""
sync_linkedin_historical.py — v0.5.0
Combina:
  (a) xls "Linkedin followers_maio 2025 ate maio 2026.xls" (diário + demografia)
  (b) 13 reports mensais PPTX (totais mensais Jan/25→Abr/26)
→ Gera public/api/linkedin-historical.json com:
  - série mensal 16 meses
  - série diária 10 meses (Aug/25-May/26)
  - eventos (tática elefante)
  - top dias sem evento
  - distribuição demográfica (5 dimensões)

Uso:
  python scripts/sync/sync_linkedin_historical.py
"""
import json, os, re, sys
from pathlib import Path
import pandas as pd

XLS_CANDIDATES = [
    # Pasta ROADMAP (Bruna atualiza aqui — fonte viva mais recente)
    r"C:/Users/Ruds/Desktop/ROADMAP MKT OFFICE JUN 2026/03 LinkedIn Boost/Linkedin followers_maio 2025 ate maio 2026.xls",
    r"C:/Users/rudac/Desktop/ROADMAP MKT OFFICE JUN 2026/03 LinkedIn Boost/Linkedin followers_maio 2025 ate maio 2026.xls",
    # Fallbacks
    r"C:/Users/Ruds/Desktop/Linkedin followers_maio 2025 ate maio 2026.xls",
    r"C:/epiuse-mkt-office/vault/00-contexto/metas/linkedin-followers-mai2025-mai2026.xls",
    r"G:/Meu Drive/Claude MKT EUBR/vault/00-contexto/metas/linkedin-followers-mai2025-mai2026.xls",
]

REPORTS_BASE = r"C:/Users/Ruds/OneDrive - EPI USE BRASIL SERVIÇOS EM SISTEMAS LTDA/MARKETING/Reports/Relatorio MKT"

OUT = Path(__file__).resolve().parents[2] / "public" / "api" / "linkedin-historical.json"


def find_xls():
    for p in XLS_CANDIDATES:
        if os.path.exists(p): return p
    return None


def extract_pptx_text(fp):
    try:
        from pptx import Presentation
    except ImportError:
        print("[linkedin] AVISO: python-pptx não instalado — pulando reports")
        return ""
    prs = Presentation(fp)
    out = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    line = ''.join(r.text for r in p.runs).strip()
                    if line: out.append(line)
    return ' '.join(out)


def parse_reports():
    """Retorna lista de {mes:'2025-01', total_seguidores, novos, newsletter, posts, impressoes, site_usuarios}"""
    results = []
    if not os.path.exists(REPORTS_BASE):
        print(f"[linkedin] AVISO: pasta reports não encontrada: {REPORTS_BASE}")
        return results

    MESES = {'janeiro':'01','fevereiro':'02','março':'03','abril':'04','maio':'05','junho':'06',
             'julho':'07','agosto':'08','setembro':'09','outubro':'10','novembro':'11','dezembro':'12'}

    for y in ['2025','2026']:
        p = f"{REPORTS_BASE}/{y}"
        if not os.path.exists(p): continue
        for f in sorted(os.listdir(p)):
            if not f.endswith('.pptx'): continue
            fp = f"{p}/{f}"
            # Extrai mês do filename: "01 - EPI-USE _ Marketing 2025 - Janeiro.pptx"
            m = re.search(r'-\s*(\w+)\.pptx$', f)
            if not m: continue
            mes_nome = m.group(1).lower()
            mes_num = MESES.get(mes_nome)
            if not mes_num: continue
            yyyymm = f"{y}-{mes_num}"
            txt = extract_pptx_text(fp)
            seg = re.search(r'([\d\.]+)\s*seguidores totais', txt)
            novos = re.search(r'([\d\.]+)\s*Novos Seguidores', txt)
            newsl = re.search(r'([\d\.]+)\s*Assinantes da Newsletter', txt)
            posts = re.search(r'(\d+)\s*posts\s*\w+', txt)
            impr = re.search(r'([\d\.]+)\s*impressões em', txt)
            results.append({
                "mes": yyyymm,
                "total_seguidores": int(seg.group(1).replace('.','')) if seg else None,
                "novos": int(novos.group(1).replace('.','')) if novos else None,
                "newsletter": int(newsl.group(1).replace('.','')) if newsl else None,
                "posts_mes": int(posts.group(1)) if posts else None,
                "impressoes": int(impr.group(1).replace('.','')) if impr else None,
                "fonte": "report-pptx",
                "arquivo": f,
            })
    # Dedupe (mantém último por mês)
    seen = {}
    for r in results: seen[r['mes']] = r
    return sorted(seen.values(), key=lambda x: x['mes'])


def main():
    fp = find_xls()
    if not fp:
        print("ERRO: xls não encontrado.")
        sys.exit(1)
    print(f"[linkedin] lendo {fp}")

    # ── Série diária + eventos
    df = pd.read_excel(fp, sheet_name='Novos seguidores')
    df['Data'] = pd.to_datetime(df['Data'], errors='coerce')
    df = df.sort_values('Data').dropna(subset=['Data'])

    col_org = [c for c in df.columns if 'rg' in c and 'nico' in c]
    col_conv = [c for c in df.columns if 'onvidados' in c]

    daily = []
    for _, r in df.iterrows():
        daily.append({
            "data": r['Data'].strftime('%Y-%m-%d'),
            "patrocinados": int(r.get('Seguidores patrocinados', 0)),
            "organicos": int(r.get(col_org[0], 0)) if col_org else 0,
            "convidados": int(r.get(col_conv[0], 0)) if col_conv else 0,
            "total": int(r.get('Total de seguidores', 0)),
            "evento": r.get('Evento') if pd.notna(r.get('Evento')) else None,
        })

    # Eventos
    eventos = [d for d in daily if d['evento']]
    eventos_sorted = sorted(eventos, key=lambda x: x['total'], reverse=True)

    # Mensal a partir do diário
    df['ym'] = df['Data'].dt.to_period('M').astype(str)
    mensal_diario = df.groupby('ym').agg(
        novos_diario=('Total de seguidores', 'sum'),
        dias_com_ganho=('Total de seguidores', lambda x: (x > 0).sum()),
    ).reset_index().rename(columns={'ym':'mes'}).to_dict('records')

    # ── Demografia
    demografia = {}
    for sheet in ['Localidade', 'Função', 'Nível de experiência', 'Setor', 'Tamanho da empresa']:
        try:
            d = pd.read_excel(fp, sheet_name=sheet)
            col_label = d.columns[0]
            col_val = d.columns[1]
            demografia[sheet] = [{"label": str(r[col_label]), "value": int(r[col_val])}
                                  for _, r in d.iterrows() if pd.notna(r[col_label])]
        except Exception as e:
            print(f"[linkedin] aviso aba {sheet}: {e}")

    # ── Série mensal dos reports (15-16 meses)
    print("[linkedin] processando reports PPTX...")
    mensal_reports = parse_reports()

    # ── Merge mensal: prefere report (tem total cumulativo), fallback pro diário
    mes_map = {r['mes']: r for r in mensal_reports}
    for m in mensal_diario:
        if m['mes'] not in mes_map:
            mes_map[m['mes']] = {
                "mes": m['mes'],
                "total_seguidores": None,
                "novos": m['novos_diario'],
                "newsletter": None,
                "posts_mes": None,
                "impressoes": None,
                "fonte": "xls-diario",
            }
        else:
            mes_map[m['mes']]['novos_diario'] = m['novos_diario']
            mes_map[m['mes']]['dias_com_ganho'] = m['dias_com_ganho']

    # ── MERGE DEFENSIVO (Regra 7 — nao regredir dado bom existente) ──────────────
    # Preserva do JSON anterior: (1) entradas fonte "manual (Ruda...)" — ex: 10640
    # cravado a mao; (2) campos cumulativos (total_seguidores/newsletter/posts/
    # impressoes) que o XLS diario nao traz, quando a fonte nova nao os tiver.
    if OUT.exists():
        try:
            prev = json.loads(OUT.read_text(encoding="utf-8"))
            prev_map = {e['mes']: e for e in prev.get('serie_mensal', [])}
            CUMUL = ['total_seguidores', 'newsletter', 'posts_mes', 'impressoes']
            for mes, old in prev_map.items():
                fonte_old = str(old.get('fonte', ''))
                is_manual = fonte_old.startswith('manual')
                new = mes_map.get(mes)
                if new is None:
                    # mes existia no JSON mas nao na fonte nova -> preserva integral
                    mes_map[mes] = old
                    continue
                # manual vence reports/diario em campos cumulativos preenchidos
                if is_manual:
                    for k in CUMUL:
                        if old.get(k) is not None:
                            new[k] = old[k]
                    if old.get('novos') is not None and new.get('novos') is None:
                        new['novos'] = old['novos']
                    # marca proveniencia combinada
                    new['fonte'] = fonte_old + ' + xls' if 'xls' not in fonte_old else fonte_old
                else:
                    # report/diario antigo: so preenche cumulativo se o novo estiver vazio
                    for k in CUMUL:
                        if new.get(k) is None and old.get(k) is not None:
                            new[k] = old[k]
        except Exception as e:
            print(f"[linkedin] aviso merge JSON anterior: {e}")

    serie_mensal = sorted(mes_map.values(), key=lambda x: x['mes'])

    total_xls = sum(d['total'] for d in daily)
    total_eventos = sum(e['total'] for e in eventos)

    payload = {
        "fonte_diario": fp,
        "fonte_reports": REPORTS_BASE,
        "gerado_em": pd.Timestamp.now().isoformat(),
        "resumo": {
            "periodo_diario": f"{daily[0]['data']} a {daily[-1]['data']}",
            "total_ganho_periodo": total_xls,
            "dias_com_ganho": sum(1 for d in daily if d['total'] > 0),
            "dias_total": len(daily),
            "media_diaria": round(total_xls / len(daily), 2),
            "total_via_eventos": total_eventos,
            "pct_eventos": round(100 * total_eventos / total_xls, 1) if total_xls else 0,
            "ultimo_total_reportado": mensal_reports[-1]['total_seguidores'] if mensal_reports else None,
        },
        "serie_mensal": serie_mensal,
        "serie_diaria": daily,
        "eventos": eventos_sorted,
        "top10_dias_sem_evento": sorted([d for d in daily if not d['evento']], key=lambda x: x['total'], reverse=True)[:10],
        "demografia": demografia,
    }

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"[linkedin] OK · {len(daily)} dias · {len(serie_mensal)} meses · {len(eventos)} eventos · {OUT.stat().st_size//1024}KB")
    print(f"[linkedin] resumo: {payload['resumo']}")


if __name__ == "__main__":
    main()
