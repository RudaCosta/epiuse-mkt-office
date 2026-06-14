#!/usr/bin/env python3
"""
gerar_pptx.py — v0.6.0
Gera PPTX mensal de Marketing EPI-USE Brasil replicando estrutura dos 13 reports históricos.

Lê: /api/relatorio/snapshot?mes=YYYY-MM (Office rodando local)
Saída: PPTX em OneDrive Marketing Reports OU caminho custom

Uso:
  python scripts/relatorio/gerar_pptx.py --mes 2026-05
  python scripts/relatorio/gerar_pptx.py --mes 2026-05 --output "C:/path/file.pptx"
"""
import argparse, json, urllib.request, urllib.error
from pathlib import Path
from datetime import datetime
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERRO: python-pptx não instalado. Rode: pip install python-pptx")
    sys.exit(1)

ONEDRIVE_BASE = r"C:/Users/Ruds/OneDrive - EPI USE BRASIL SERVIÇOS EM SISTEMAS LTDA/MARKETING/Reports/Relatorio MKT"

MESES_PT = {
    "01": "Janeiro", "02": "Fevereiro", "03": "Março", "04": "Abril",
    "05": "Maio", "06": "Junho", "07": "Julho", "08": "Agosto",
    "09": "Setembro", "10": "Outubro", "11": "Novembro", "12": "Dezembro",
}

# Paleta EPI-USE (aproximadamente igual ao branding)
COR_PRIMARY = RGBColor(0x13, 0x1B, 0x41)   # navy ERP.ngo
COR_SECONDARY = RGBColor(0x00, 0x66, 0xB2)  # mid blue
COR_ACCENT = RGBColor(0xFB, 0xBF, 0x24)     # amarelo destaque
COR_GREEN = RGBColor(0x10, 0xB9, 0x81)
COR_RED = RGBColor(0xEF, 0x44, 0x44)
COR_TEXT = RGBColor(0x1E, 0x29, 0x3B)
COR_MUTED = RGBColor(0x64, 0x74, 0x8B)
COR_WHITE = RGBColor(255, 255, 255)


def get_layout(prs, name, fallback_idx=6):
    """Retorna um layout do template buscando por nome, ou fallback por index."""
    for l in prs.slide_layouts:
        if l.name == name:
            return l
    if fallback_idx < len(prs.slide_layouts):
        return prs.slide_layouts[fallback_idx]
    return prs.slide_layouts[0]


def fetch_snapshot(mes, base_url="http://localhost:3000"):
    url = f"{base_url}/api/relatorio/snapshot?mes={mes}"
    try:
        with urllib.request.urlopen(url, timeout=30) as r:
            return json.loads(r.read().decode("utf-8"))
    except urllib.error.URLError as e:
        print(f"ERRO: não conseguiu acessar {url}. Office tá rodando? ({e})")
        sys.exit(1)


def add_titulo(slide, texto, top=Inches(0.3), size=24, bold=True):
    tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.5), Inches(0.7))
    tf = tb.text_frame
    tf.text = texto
    tf.paragraphs[0].runs[0].font.size = Pt(size)
    tf.paragraphs[0].runs[0].font.bold = bold
    tf.paragraphs[0].runs[0].font.color.rgb = COR_PRIMARY


def add_kpi_card(slide, x, y, w, h, label, valor, sub=None, cor_val=COR_SECONDARY):
    """Card retangular com label + valor grande + sub opcional"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(0.5)

    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = label
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = COR_MUTED
    p.runs[0].font.bold = True

    tb2 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.4), w - Inches(0.3), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = str(valor)
    p2.runs[0].font.size = Pt(22)
    p2.runs[0].font.bold = True
    p2.runs[0].font.color.rgb = cor_val
    p2.runs[0].font.name = "Consolas"

    if sub:
        tb3 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.0), w - Inches(0.3), Inches(0.3))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = sub
        p3.runs[0].font.size = Pt(9)
        p3.runs[0].font.color.rgb = COR_MUTED


def slide_capa(prs, mes):
    layout = get_layout(prs, "Title-slide_Elephant", 6)
    s = prs.slides.add_slide(layout)
    yyyy, mm = mes.split("-")
    nome_mes = MESES_PT[mm]

    if layout.name == "Title-slide_Elephant":
        # Se for o layout oficial, preenche os placeholders 14 e 15 com texto branco e alinhado
        try:
            ph_title = s.placeholders[14]
            ph_title.text = "Monthly Report — Marketing"
        except KeyError:
            pass
        try:
            ph_sub = s.placeholders[15]
            ph_sub.text = f"{nome_mes} de {yyyy}  ·  EPI-USE Brasil"
        except KeyError:
            pass
    else:
        # Fallback original
        add_titulo(s, "Monthly Report — Marketing", top=Inches(2.5), size=36)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.5), Inches(0.7))
        tb.text_frame.text = f"{nome_mes} de {yyyy}"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_SECONDARY

        tb2 = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.5), Inches(0.4))
        tb2.text_frame.text = f"EPI-USE Brasil · Marketing & RevOps · Gerado em {datetime.now().strftime('%d/%m/%Y %H:%M')}"
        tb2.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        tb2.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_MUTED


def slide_agenda(prs):
    layout = get_layout(prs, "Blue BG", 6)
    s = prs.slides.add_slide(layout)
    is_blue_bg = "Blue BG" in layout.name

    itens = [
        ("1.", "KPIs Digitais"),
        ("2.", "LinkedIn — Detalhe"),
        ("3.", "Conteúdo"),
        ("4.", "Eventos (Tática Elefante)"),
        ("5.", "EPI-USE Voices"),
        ("6.", "SDR / Pipeline"),
        ("7.", "Next Steps"),
    ]

    if is_blue_bg:
        # Título no meio-esquerda em branco
        tb = s.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COR_WHITE

        y = Inches(2.5)
        for num, label in itens:
            tb_item = s.shapes.add_textbox(Inches(1.0), y, Inches(11.3), Inches(0.4))
            p_item = tb_item.text_frame.paragraphs[0]
            p_item.text = f"{num}  {label}"
            p_item.font.size = Pt(18)
            p_item.font.color.rgb = COR_WHITE
            y += Inches(0.5)
    else:
        # Comportamento original
        add_titulo(s, "Agenda")
        y = Inches(1.5)
        for num, label in itens:
            tb = s.shapes.add_textbox(Inches(2), y, Inches(10), Inches(0.5))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{num}  {label}"
            p.runs[0].font.size = Pt(18)
            p.runs[0].font.color.rgb = COR_PRIMARY
            y += Inches(0.5)



def slide_kpis_digitais(prs, snap):
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "1. KPIs Digitais")
    li = snap["linkedin"]
    site = snap.get("site")
    insta = snap.get("instagram")
    em = snap.get("email")

    # Preenche Site
    site_kpis = []
    if site and site.get("usuarios") is not None:
        site_kpis.append(("Usuários", f"{site['usuarios']}"))
        site_kpis.append(("Visualizações", f"{site['visualizacoes']}"))
        site_kpis.append(("Tempo engajamento", f"{site['duracao_sessao_s']}s" if site.get("duracao_sessao_s") else "—"))
    else:
        site_kpis.extend([("Usuários", "—"), ("Visualizações", "—"), ("Tempo engajamento", "—"), ("⏳ Aguarda GA4 API", "")])

    # Preenche Instagram
    insta_kpis = []
    if insta and insta.get("seguidores_novos") is not None:
        insta_kpis.append(("Novos seguidores", f"+{insta['seguidores_novos']}"))
        insta_kpis.append(("Alcance", f"{insta['alcance']}"))
    else:
        insta_kpis.extend([("Seguidores", "—"), ("Alcance", "—"), ("⏳ Aguarda Graph API", "")])

    # Preenche E-mail
    email_kpis = []
    if em and em.get("taxa_abertura") is not None:
        email_kpis.append(("Taxa abertura", f"{em['taxa_abertura']}%"))
        email_kpis.append(("Taxa cliques", f"{em['taxa_cliques']}%"))
        email_kpis.append(("Leads (conversões)", f"{em['leads']}"))
    else:
        email_kpis.extend([("Taxa abertura", "—"), ("Taxa cliques", "—"), ("Leads", "—"), ("⏳ Aguarda Token RD", "")])

    # 4 colunas: Site · LinkedIn · Instagram · E-mail
    col_w = Inches(2.85)
    col_h = Inches(4.5)
    x = Inches(0.5)
    y = Inches(1.5)
    canais = [
        ("🌐 Site", site_kpis),
        ("💼 LinkedIn", [
            ("Seguidores", f"{li.get('total_atual') or '—'}"),
            ("Novos no mês", f"+{li.get('novos') or '—'}"),
            ("Impressões", f"{li.get('impressoes') or '—'}"),
            ("Engajamento", f"{li.get('engajamento')}%" if li.get("engajamento") is not None else "—"),
        ]),
        ("📷 Instagram", insta_kpis),
        ("✉️ E-mail (RD)", email_kpis),
    ]
    for nome, kpis in canais:
        # Header do card
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, col_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

        tb = s.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = nome
        p.runs[0].font.size = Pt(14); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = COR_PRIMARY

        yy = y + Inches(0.6)
        for lbl, val in kpis:
            tb2 = s.shapes.add_textbox(x + Inches(0.15), yy, col_w - Inches(0.3), Inches(0.45))
            tf = tb2.text_frame
            tf.text = lbl
            tf.paragraphs[0].runs[0].font.size = Pt(9); tf.paragraphs[0].runs[0].font.color.rgb = COR_MUTED
            if val:
                pp = tf.add_paragraph()
                pp.text = str(val)
                pp.runs[0].font.size = Pt(13); pp.runs[0].font.bold = True
                pp.runs[0].font.color.rgb = COR_SECONDARY
                pp.runs[0].font.name = "Consolas"
            yy += Inches(0.9)
        x += col_w + Inches(0.05)


def slide_linkedin(prs, snap):
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "2. LinkedIn — Detalhe")
    li = snap["linkedin"]
    add_kpi_card(s, Inches(0.5), Inches(1.5), Inches(3.5), Inches(1.5),
                 "Seguidores Totais", f"{li.get('total_atual') or '—'}",
                 f"+{li.get('novos') or '—'} novos no mês", COR_SECONDARY)
    add_kpi_card(s, Inches(4.2), Inches(1.5), Inches(3.5), Inches(1.5),
                 "Newsletter", f"{li.get('newsletter') or '—'}", "assinantes", COR_PRIMARY)
    add_kpi_card(s, Inches(7.9), Inches(1.5), Inches(3.5), Inches(1.5),
                 "Impressões", f"{li.get('impressoes') or '—'}", f"{li.get('posts_mes') or '—'} posts", COR_GREEN)

    # Demografia (top 3 de cada dimensão)
    demo = li.get("demografia", {})
    add_titulo(s, "Demografia dos Seguidores", top=Inches(3.5), size=14)
    yy = Inches(4.0)
    dims = [("Localidade", "📍"), ("Função", "💼"), ("Nível de experiência", "🎓"), ("Setor", "🏢")]
    for i, (k, ico) in enumerate(dims):
        items = demo.get(k, [])[:3]
        x = Inches(0.5 + i * 3.1)
        tb = s.shapes.add_textbox(x, yy, Inches(3), Inches(0.3))
        tb.text_frame.text = f"{ico} {k}"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(10); tb.text_frame.paragraphs[0].runs[0].font.bold = True
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_MUTED
        for j, it in enumerate(items):
            tb2 = s.shapes.add_textbox(x, yy + Inches(0.35 + j*0.3), Inches(3), Inches(0.3))
            tb2.text_frame.text = f"{it['label'][:30]} · {it['value']:,}"
            tb2.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
            tb2.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_TEXT


def slide_eventos(prs, snap):
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "4. Eventos — Tática Elefante 🐘")
    tat = snap["linkedin"].get("tatica_elefante", {})
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
    tb.text_frame.text = f"{tat.get('eventos_no_periodo', 0)} eventos · +{tat.get('seguidores_via_eventos', 0)} seguidores ({tat.get('pct_eventos', 0)}% do total orgânico)"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(14); tb.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_ACCENT

    # Próximos
    add_titulo(s, "📅 Próximos Eventos", top=Inches(3.5), size=14)
    eventos = snap.get("eventos_proximos", [])
    yy = Inches(4.0)
    for e in eventos[:5]:
        nome = e.get("nome") or e.get("titulo") or e.get("evento", "—")
        data = e.get("data_inicio") or e.get("data", "")
        tb = s.shapes.add_textbox(Inches(0.5), yy, Inches(12), Inches(0.4))
        tb.text_frame.text = f"• {nome}   ({data})"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        yy += Inches(0.45)


def slide_voices(prs, snap):
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "5. EPI-USE Voices")
    v = snap.get("voices", {})
    add_kpi_card(s, Inches(0.5), Inches(1.5), Inches(3.5), Inches(1.5),
                 "Voices Ativos", f"{v.get('ativos', 0)}", f"de {v.get('total', 0)} total")
    yy = Inches(3.3)
    for voice in v.get("lista", [])[:6]:
        tb = s.shapes.add_textbox(Inches(0.5), yy, Inches(12), Inches(0.4))
        tb.text_frame.text = f"• {voice['nome']} — {voice.get('area', '—')} · SSI {voice.get('ssi', '?')} · {voice.get('seg') or '?'} seguidores"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        yy += Inches(0.45)


def slide_cases(prs, snap):
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "7. Cases & CS Highlights")
    c = snap["cases"]
    cards = [
        ("Publicado", c.get("publicado", 0), COR_GREEN),
        ("Em Edição", c.get("em_edicao", 0), COR_ACCENT),
        ("Em Negociação", c.get("negociacao", 0), COR_SECONDARY),
        ("Declinado", c.get("declinado", 0), COR_MUTED),
    ]
    x = Inches(0.5)
    for lbl, val, cor in cards:
        add_kpi_card(s, x, Inches(1.5), Inches(2.9), Inches(1.5), lbl, str(val), None, cor)
        x += Inches(3.05)


def slide_alertas(prs, snap):
    alertas = snap.get("alertas", [])
    if not alertas: return
    layout = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout)
    add_titulo(s, "⚠️ Alertas Automáticos")
    yy = Inches(1.5)
    for a in alertas:
        cor = COR_ACCENT if a["tipo"] == "warn" else COR_SECONDARY
        tb = s.shapes.add_textbox(Inches(0.5), yy, Inches(12), Inches(0.5))
        tb.text_frame.text = f"{'⚠️' if a['tipo']=='warn' else 'ℹ️'} {a['msg']}"
        tb.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        tb.text_frame.paragraphs[0].runs[0].font.color.rgb = cor
        yy += Inches(0.55)


def slide_fim(prs):
    layout = get_layout(prs, "End-slide_Elephant", 6)
    if layout.name == "End-slide_Elephant":
        s = prs.slides.add_slide(layout)
        try:
            ph = s.placeholders[12]
            ph.text = "Obrigado!"
        except KeyError:
            pass


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--mes", required=True, help="YYYY-MM (ex: 2026-05)")
    ap.add_argument("--output", help="Path custom do PPTX")
    ap.add_argument("--base-url", default="http://localhost:3000", help="URL do Office")
    ap.add_argument("--template", help="Caminho do template PPTX da EPI-USE")
    args = ap.parse_args()

    print(f"[relatorio] pegando snapshot {args.mes}...")
    snap = fetch_snapshot(args.mes, args.base_url)

    print(f"[relatorio] montando PPTX...")
    
    template_path = None
    if args.template:
        template_path = Path(args.template)
    else:
        template_path = Path(ONEDRIVE_BASE) / "EPI-USE_PPT_Template (Effective Jan 2026).pptx"

    if template_path.exists():
        print(f"[relatorio] Carregando template oficial: {template_path}")
        prs = Presentation(str(template_path))
        # Deletar todos os slides originais do template
        for i in range(len(prs.slides) - 1, -1, -1):
            slide_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[i]
    else:
        print("[relatorio] Template oficial não encontrado. Usando apresentação vazia padrão.")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slide_capa(prs, args.mes)
    slide_agenda(prs)
    slide_kpis_digitais(prs, snap)
    slide_linkedin(prs, snap)
    
    # Conteúdo (placeholder por enquanto)
    layout_content = get_layout(prs, "Content-slide_white-bg-blank", 6)
    s = prs.slides.add_slide(layout_content)
    add_titulo(s, "3. Conteúdo do Mês")
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(0.5))
    tb.text_frame.text = "[Aguardando tracking de Conteúdo em /api — temas trabalhados, taxa engajamento, destaques]"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(13); tb.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_MUTED

    slide_eventos(prs, snap)
    slide_voices(prs, snap)

    # SDR
    s_sdr = prs.slides.add_slide(layout_content)
    add_titulo(s_sdr, "6. SDR / Pipeline")
    tb_sdr = s_sdr.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(0.5))
    tb_sdr.text_frame.text = "[Aguardando integração Apollo MCP — Onda 5 v0.5.0]"
    tb_sdr.text_frame.paragraphs[0].runs[0].font.size = Pt(13); tb_sdr.text_frame.paragraphs[0].runs[0].font.color.rgb = COR_MUTED

    slide_cases(prs, snap)
    slide_alertas(prs, snap)
    slide_fim(prs)

    # Output path
    if args.output:
        out = Path(args.output)
    else:
        yyyy, mm = args.mes.split("-")
        nome_mes = MESES_PT[mm]
        out = Path(ONEDRIVE_BASE) / yyyy / f"{mm} - EPI-USE _ Marketing {yyyy} - {nome_mes} (auto).pptx"

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"[relatorio] OK")
    print(f"           PPTX: {out}")
    print(f"           Dashboard: {args.base_url}/relatorio?mes={args.mes}")
    print(f"           Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
